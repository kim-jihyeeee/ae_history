import streamlit as st
import pandas as pd
from wordcloud import WordCloud
import matplotlib.pyplot as plt
from io import BytesIO
import datetime, os, requests, re
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches

# 1. 기본 설정 및 폰트 로드
st.set_page_config(page_title="AE Total Solution v9.0", layout="wide")

@st.cache_data
def load_font():
    try:
        url = "https://github.com/google/fonts/raw/main/ofl/nanumgothic/NanumGothic-Bold.ttf"
        res = requests.get(url)
        with open("nanum_font.ttf", "wb") as f: f.write(res.content)
        return "nanum_font.ttf"
    except: return None

FONT_PATH = load_font()

# 세션 초기화
if 'client_db' not in st.session_state: st.session_state.client_db = pd.DataFrame()
if 'history_db' not in st.session_state: st.session_state.history_db = pd.DataFrame(columns=['날짜', '광고주명', '소통내용', '핵심키워드'])

# 2. UI 스타일 (메뉴 간격 및 버튼 디자인)
st.markdown("""
    <style>
    header[data-testid="stHeader"] { visibility: hidden; }
    .stButton>button { width: 100%; border-radius: 8px; background-color: #FFB300; color: white; font-weight: bold; height: 3em; }
    .menu-header { font-size: 1.1em; font-weight: bold; color: #FFB300; margin-top: 35px; border-bottom: 2px solid #eee; padding-bottom: 5px; }
    .spacer { margin-bottom: 50px; }
    div[role="slider"] { background-color: #FF4B4B !important; }
    </style>
""", unsafe_allow_html=True)

# 3. 사이드바 메뉴 (내부/외부 확실한 구분)
st.sidebar.title("🚀 AE Total Tool v9.0")
st.sidebar.markdown('<p class="menu-header">📋 내부 히스토리 관리</p>', unsafe_allow_html=True)
m_int = st.sidebar.radio("항목", ["광고주 DB 관리", "관리 이력 입력", "디지털 리포트(내부)"], label_visibility="collapsed")

st.sidebar.markdown('<div class="spacer"></div>', unsafe_allow_html=True)

st.sidebar.markdown('<p class="menu-header">📊 외부 시장 분석</p>', unsafe_allow_html=True)
m_ext = st.sidebar.checkbox("📊 Trend Radar(듀얼 모드)")

menu = "📊 Trend Radar(외부)" if m_ext else m_int

# --- [메뉴 1: 광고주 DB 관리] ---
if menu == "광고주 DB 관리":
    st.header("📂 데이터 로드 및 관리")
    c1, c2 = st.columns(2)
    with c1: 
        up_c = st.file_uploader("광고주 리스트 (엑셀/CSV)", type=['xlsx', 'csv'], key="c_up")
        if up_c:
            df = pd.read_csv(up_c) if up_c.name.endswith('.csv') else pd.read_excel(up_c)
            df.rename(columns={df.columns[0]: '광고주명'}, inplace=True)
            st.session_state.client_db = df
            st.success("✅ 광고주 리스트 로드 완료!")
    with c2:
        up_h = st.file_uploader("히스토리 백업 (엑셀)", type=['xlsx'], key="h_up")
        if up_h:
            st.session_state.history_db = pd.read_excel(up_h)
            st.success("✅ 히스토리 복구 완료!")

# --- [메뉴 2: 관리 이력 입력] ---
elif menu == "관리 이력 입력":
    st.header("✍️ 소통 이력 기록 및 수정")
    if st.session_state.client_db.empty:
        st.warning("먼저 '광고주 DB 관리' 메뉴에서 리스트를 등록하세요.")
    else:
        q = st.text_input("🔍 업체명 검색")
        all_c = sorted(st.session_state.client_db['광고주명'].dropna().unique())
        filtered_c = [c for c in all_c if q.lower() in str(c).lower()]
        
        with st.form("history_form"):
            c1, c2 = st.columns(2)
            with c1:
                sel_c = st.selectbox(f"광고주 선택 ({len(filtered_c)}건)", filtered_c)
                dt = st.date_input("날짜", datetime.date.today())
            with c2: tags = st.text_input("핵심 키워드 (쉼표 구분)")
            txt = st.text_area("상세 소통 내용", height=150)
            if st.form_submit_button("히스토리 저장"):
                new_row = pd.DataFrame([[dt, sel_c, txt, tags]], columns=['날짜', '광고주명', '소통내용', '핵심키워드'])
                st.session_state.history_db = pd.concat([st.session_state.history_db, new_row], ignore_index=True)
                st.rerun()
        st.divider()
        st.subheader("🛠️ 전체 데이터 편집기")
        st.data_editor(st.session_state.history_db, use_container_width=True)

# --- [메뉴 3: 디지털 리포트 (내부)] ---
elif menu == "디지털 리포트(내부)":
    st.header("📊 내부 소통 키워드 분석")
    if st.session_state.history_db.empty:
        st.info("기록된 데이터가 없습니다.")
    else:
        target = st.selectbox("광고주 선택", sorted(st.session_state.history_db['광고주명'].unique()))
        f = st.session_state.history_db[st.session_state.history_db['광고주명'] == target]
        if not f.empty:
            words = (f['핵심키워드'].fillna('').str.cat(sep=' ') + " ") * 3 + f['소통내용'].fillna('').str.cat(sep=' ')
            wc = WordCloud(font_path=FONT_PATH, width=800, height=400, background_color='white').generate(words)
            fig, ax = plt.subplots(); ax.imshow(wc); ax.axis('off'); st.pyplot(fig)

# --- [메뉴 4: 🌟 Trend Radar 듀얼 모드 (외부)] ---
elif menu == "📊 Trend Radar(외부)":
    st.header("🌐 Trend Radar 듀얼 분석 (뉴스 & 검색)")
    tab_news, tab_search = st.tabs(["📰 뉴스 에디션 (기사 분석)", "🔍 검색 에디션 (포털 분석)"])

    # 1. 뉴스 에디션 (RSS 방식 - 보안 차단 없음)
    with tab_news:
        st.subheader("최신 뉴스 기반 시장 트렌드")
        cn1, cn2 = st.columns([3, 1])
        with cn1: kw_n = st.text_input("뉴스 분석 키워드", placeholder="예: 무라벨 생수 트렌드")
        with cn2: prd_n = st.selectbox("수집 기간", ["최근 7일", "최근 30일", "최근 60일"], key="news_prd")
        
        if st.button("🔍 뉴스 트렌드 분석 시작"):
            with st.spinner("뉴스를 정밀 분석 중입니다..."):
                days = int(re.findall(r'\d+', prd_n)[0])
                limit = datetime.datetime.now() - datetime.timedelta(days=days)
                rss_url = f"https://news.google.com/rss/search?q={kw_n}&hl=ko&gl=KR&ceid=KR:ko"
                res = requests.get(rss_url)
                soup = BeautifulSoup(res.text, 'xml')
                items = soup.find_all('item')
                
                filtered = []
                for i in items:
                    p_date = datetime.datetime.strptime(i.pubDate.get_text(), '%a, %d %b %Y %H:%M:%S %Z')
                    if p_date >= limit:
                        clean_t = re.split(r' - | \| ', i.title.get_text())[0]
                        filtered.append(clean_t)
                
                if not filtered: st.error("해당 기간 뉴스 없음")
                else:
                    news_txt = " ".join(filtered)
                    # 매체명 제거 필터
                    for m in ["동아일보", "서울경제", "헬스조선", "하이닥", "연합뉴스", "매일경제", "한국경제", "네이버", "뉴스", "daum", "net"]:
                        news_txt = news_txt.replace(m, "")
                    
                    wc_n = WordCloud(font_path=FONT_PATH, width=900, height=500, background_color='white', colormap='cool', regexp=r"[\w\xA1-\xFE]+").generate(news_txt)
                    fig_n, ax_n = plt.subplots(figsize=(10, 6)); ax_n.imshow(wc_n, interpolation='bilinear'); ax_n.axis('off'); st.pyplot(fig_n)
                    
                    st.divider()
                    dn1, dn2 = st.columns(2)
                    buf_n = BytesIO(); fig_n.savefig(buf_n, format="png", dpi=300)
                    dn1.download_button("📥 뉴스 이미지 저장", data=buf_n.getvalue(), file_name=f"News_{kw_n}.png")
                    
                    try:
                        prs_n = Presentation(); slide_n = prs_n.slides.add_slide(prs_n.slide_layouts[6])
                        slide_n.shapes.add_picture(buf_n, Inches(0.5), Inches(1), width=Inches(9))
                        ppt_n = BytesIO(); prs_n.save(ppt_n)
                        dn2.download_button("📊 뉴스 PPTX 저장", data=ppt_n.getvalue(), file_name=f"News_{kw_n}.pptx")
                    except: pass

    # 2. 검색 에디션 (포털 정밀 수집 방식 - 보안 우회)
    with tab_search:
        st.subheader("포털 실시간 검색 결과 분석")
        kw_s = st.text_input("검색 분석 키워드", placeholder="예: 단백질 쉐이크 추천, 슬개골 탈구 영양제")
        
        if st.button("🔍 소비자 검색 분석 시작"):
            with st.spinner("포털 데이터를 안전하게 수집 중입니다..."):
                headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"}
                search_all = ""
                
                # 네이버 & 다음 교차 수집 (구글 보안 우회)
                sources = [f"https://search.naver.com/search.naver?query={kw_s}", f"https://search.daum.net/search?q={kw_s}"]
                for url in sources:
                    try:
                        res = requests.get(url, headers=headers, timeout=10)
                        soup = BeautifulSoup(res.text, 'html.parser')
                        # 제목과 설명 텍스트만 추출
                        search_all += " ".join([v.get_text() for v in soup.select('.lnk_tit, .dsc_txt, .api_txt_lines, .total_tit, .tit_main, .desc_txt')])
                    except: pass

                # 보안 문구 및 불용어 정제
                for trash in ["닫기", "열기", "이동", "클릭", "Search", "Google", "Please", "click", "accessing", "redirected"]:
                    search_all = re.sub(trash, "", search_all, flags=re.I)

                if len(search_all.strip()) < 50:
                    st.error("데이터 수집 실패. 키워드를 더 구체적으로 입력하세요.")
                else:
                    wc_s = WordCloud(font_path=FONT_PATH, width=900, height=500, background_color='white', colormap='YlOrRd', regexp=r"[\w\xA1-\xFE]+").generate(search_all)
                    fig_s, ax_s = plt.subplots(figsize=(10, 6)); ax_s.imshow(wc_s, interpolation='bilinear'); ax_s.axis('off'); st.pyplot(fig_s)
                    
                    st.divider()
                    ds1, ds2 = st.columns(2)
                    buf_s = BytesIO(); fig_s.savefig(buf_s, format="png", dpi=300)
                    ds1.download_button("📥 검색 이미지 저장", data=buf_s.getvalue(), file_name=f"Search_{kw_s}.png")
                    
                    try:
                        prs_s = Presentation(); slide_s = prs_s.slides.add_slide(prs_s.slide_layouts[6])
                        slide_s.shapes.add_picture(buf_s, Inches(0.5), Inches(1), width=Inches(9))
                        ppt_s = BytesIO(); prs_s.save(ppt_s)
                        ds2.download_button("📊 검색 PPTX 저장", data=ppt_s.getvalue(), file_name=f"Search_{kw_s}.pptx")
                    except: pass
